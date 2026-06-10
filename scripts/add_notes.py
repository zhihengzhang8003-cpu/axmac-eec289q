"""Write speaker notes to all 18 slides of final.pptx.

Usage:
    py -3.14 scripts/add_notes.py
"""
import sys
from pathlib import Path
from copy import deepcopy
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

sys.stdout.reconfigure(encoding="utf-8")

PPTX = Path(r"E:\AIDev\EEC 289Q 002 SQ 2026：Deep Learning Hardware\project\final.pptx")

NOTES = {
    0: (  # Slide 1: Title
        "Good [morning/afternoon]. I'm [name], presenting joint work with my partner. "
        "Today's talk is about AxMAC — a framework for evaluating approximate MAC designs "
        "for energy-efficient DNN inference. We'll cover the theory, simulation, and "
        "on-chip FPGA validation."
    ),
    1: (  # Slide 2: Previous Work
        "Four papers motivate this work. "
        "Han et al. showed DNN inference is both compute- and memory-intensive — "
        "justifying hardware-level optimization. "
        "Gupta et al. showed that the rounding scheme matters even at low precision. "
        "Wu et al. reviewed approximate multiplier designs, confirming energy tradeoffs. "
        "Sayadi et al. showed different layers tolerate different approximation levels — "
        "directly motivating our Contribution B."
    ),
    2: (  # Slide 3: Introduction & Motivation
        "MAC operations dominate DNN inference — often 80–90% of total compute. "
        "The key observation: a product has many bits, but low-significance bits "
        "contribute little to the final result. "
        "Truncating K of those low bits simplifies the adder tree, reducing area and power. "
        "The central question: how much error does this introduce, and can we control it?"
    ),
    3: (  # Slide 4: Architecture overview (image slide)
        "This diagram shows the five-module framework. "
        "exact_mac and approx_mac are the scalar arithmetic layer. "
        "dnn_inference lifts them to network-level evaluation. "
        "sensitivity implements Contribution B's layer profiling and allocation. "
        "pareto extracts the non-dominated design frontier across all configurations."
    ),
    4: (  # Slide 5: exact_mac.py
        "exact_mac.py is the reference baseline — every approximate result is compared against it. "
        "It implements Booth radix-4 partial products for INT "
        "and full sign/exponent/mantissa decomposition for FP. "
        "Setting K=0 and no ACA window recovers exact arithmetic exactly."
    ),
    5: (  # Slide 6: approx_mac.py
        "approx_mac.py introduces three approximation knobs: "
        "K-bit product truncation, rounding mode (trunc / round / stochastic), "
        "and ACA carry window W. "
        "K=0 and W=None reproduce the exact path. "
        "This module is the hardware primitive being evaluated throughout."
    ),
    6: (  # Slide 7: dnn_inference.py
        "dnn_inference.py propagates MAC-level error through a full network. "
        "Every product x[i,k] × w[k,j] is truncated before accumulation. "
        "tiny_mlp_forward accepts a scalar K for uniform truncation "
        "or a per-layer K list for Contribution B. "
        "This is where single-MAC errors become DNN-level output degradation."
    ),
    7: (  # Slide 8: Three Ways to Drop Low K Bits
        "Walk left to right. "
        "Trunc: drop K bits, no compensation, errors always negative, O(N) accumulation. "
        "Round: add 2^(K-1) before dropping, errors cancel, near-zero bias, one constant adder. "
        "Stochastic: add random [0, 2^K) before dropping, truly zero-mean, needs LFSR — priciest. "
        "Cost rises left to right; bias removal also improves left to right."
    ),
    8: (  # Slide 9: Contribution A
        "Core finding: trunc bias grows O(N), not O(√N). "
        "At N=4096, K=6: trunc=26,617 vs stochastic=14 — 1,900× difference. "
        "Theorem 1 formalizes this. "
        "Round's hardware overhead: +0.53 mW, zero registers — one constant adder. "
        "Transformer Corollary: d_head=64, K=4 → 512 LSB per attention score, "
        "exceeds INT8 range — silent softmax saturation."
    ),
    9: (  # Slide 10: Pareto Analysis
        "Two conflicting objectives: energy and error — can't minimize both at once. "
        "Design A dominates B if A is no worse on both axes and strictly better on one. "
        "The Pareto front is the lower-left boundary after removing all dominated designs. "
        "Every point on the front is a genuine tradeoff; anything off it is strictly worse."
    ),
    10: (  # Slide 11: Results Energy-Error
        "Three precisions fall into distinct energy bands — "
        "INT4 around 0.08–0.11 pJ, INT8 at 0.20–0.28 pJ, INT16 at 0.56–0.69 pJ. "
        "INT16 is nearly flat to K=12 — plenty of headroom. "
        "INT8 knee point is K=2: 15% energy savings, NRMSE rises from 0.148 to 0.172. "
        "Beyond K=2, INT8 error climbs sharply. "
        "INT4 is cheapest but the error curve is steep — limited precision, little tolerance."
    ),
    11: (  # Slide 12: Contribution B
        "Uniform K is wasteful — layers differ in both MAC count and sensitivity. "
        "First layer: 200,704 MACs, 95.8% of total, low sensitivity. "
        "Output layer: 640 MACs, highly sensitive — errors distort class scores directly. "
        "Method: probe each layer individually → sensitivity score → greedy allocator. "
        "Output layer always pinned at K=0. "
        "Result at B=1: 10.9× lower NRMSE than uniform K at identical energy."
    ),
    12: (  # Slide 13: RTL Module Map
        "Three levels. Bottom: mac_unit (all three rounding modes) + lfsr (64-bit, for stochastic) "
        "+ aca_adder (parameterized by window W). "
        "Middle: mac_array tiles a 1×4 grid of MAC units. "
        "Top: mlp_top drives the array through the MLP with an 8-state FSM, ~320 cycles per inference. "
        "Board wrapper mlp_top_demo adds UART output — 41-byte logit frame at 9600 baud. "
        "All modules passed ModelSim testbenches against Python golden vectors."
    ),
    13: (  # Slide 14: RTL Results
        "Synthesis on EP4CE10 at 50 MHz. "
        "Trunc: 1,685 LE, 527 regs, 8.67 mW — cheapest baseline. "
        "Round: +46 LE, +0 regs, +0.53 mW — one constant-addition corrector, no state. "
        "Stochastic: +132 LE, +64 regs, +0.72 mW — those 64 regs ARE the LFSR. "
        "Stochastic at 1,817 LE is larger than exact K=0 at 1,769 LE — RNG costs more than truncation saves. "
        "Round: near-zero bias at 35% of stochastic's logic overhead, zero register cost. "
        "K=6 truncation: 17% core power savings vs exact."
    ),
    14: (  # Slide 15: FPGA Validation
        "Five bitstream configs burned via JTAG, argmax read over UART. "
        "Test input: correctly labeled class 1 under exact computation. "
        "K=0 / K=2 / K=4 trunc and K=4 round: all output class 1 — PASS. "
        "K=6 trunc: outputs class 3 — FAIL. "
        "Theorem 1 predicts bias = 4096 × 32 = 131,072 LSB — far beyond INT8 range. "
        "K=4 round passes because round cancels systematic bias entirely."
    ),
    15: (  # Slide 16: Conclusions
        "Contribution A: trunc bias is O(N). "
        "Round fixes it with one constant adder — 0.53 mW, zero registers. "
        "Transformer Corollary: silent softmax saturation at d_head=64, K=4. "
        "Contribution B: sensitivity-driven allocation → 10.9× lower NRMSE at same budget. "
        "FPGA validation: K=6 trunc misclassifies exactly as predicted; K=4 round passes. "
        "Future: NanGate45 ASIC, ResNet-50/ImageNet, GPT-2 perplexity."
    ),
    16: (  # Slide 17: References
        "Four references cited throughout the talk."
    ),
    17: (  # Slide 18: Thank You
        "Thank you. Happy to take any questions."
    ),
}


def set_notes(slide, text: str) -> None:
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)
    first_para = tf.paragraphs[0]
    for run in first_para.runs:
        run._r.getparent().remove(run._r)
    lines = text.split("\n")
    first = True
    for line in lines:
        if first:
            first_para.text = line
            first = False
        else:
            new_p = deepcopy(first_para._p)
            for r in new_p.findall(qn("a:r")):
                new_p.remove(r)
            r_elem = etree.SubElement(new_p, qn("a:r"))
            t_elem = etree.SubElement(r_elem, qn("a:t"))
            t_elem.text = line
            tf._txBody.append(new_p)


def main():
    prs = Presentation(str(PPTX))
    total = len(prs.slides)
    print(f"Loaded {total} slides")
    for idx, text in NOTES.items():
        if idx >= total:
            print(f"  SKIP slide {idx+1}")
            continue
        set_notes(prs.slides[idx], text)
        print(f"  Slide {idx+1}: {len(text)} chars")
    prs.save(str(PPTX))
    print("Saved.")


if __name__ == "__main__":
    main()
