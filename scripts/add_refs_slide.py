import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
_PPTX = os.path.join(_ROOT, "deliverables", "AxMAC_Defense_Slides.pptx")
prs = Presentation(_PPTX)

W = prs.slide_width
H = prs.slide_height

blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# Background
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x0F, 0x11, 0x17)

WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x7A, 0x80, 0xA0)

def add_textbox(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = txb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txb

# Title
add_textbox(slide,
    l=457200, t=228600, w=11277295, h=533400,
    text='References', size=20, color=WHITE, bold=True)

# Slide number
add_textbox(slide,
    l=10972695, t=6400800, w=1219200, h=304800,
    text='12 / 12', size=9, color=GRAY, bold=False, align=PP_ALIGN.RIGHT)

# Reference list
refs = [
    "[1]  H. Mahdiani et al., \"Bio-inspired imprecise computational blocks for efficient VLSI implementation,\" IEEE TCAS-I, 2010.",
    "[2]  M. J. Schulte & E. E. Swartzlander, \"Truncated multiplication with correction constant,\" VLSI Signal Processing VI, 1993.",
    "[3]  S. Gupta et al., \"Deep learning with limited numerical precision,\" ICML, 2015.",
    "[4]  A. B. Kahng & S. Kang, \"Accuracy-configurable adder for approximate arithmetic designs,\" DAC, 2012.",
    "[5]  S. Hashemi et al., \"DRUM: A dynamic range unbiased multiplier for approximate applications,\" ICCAD, 2015.",
    "[6]  M. Horowitz, \"Computing's energy problem (and what we can do about it),\" ISSCC 2014 keynote.",
    "[7]  P. Micikevicius et al., \"FP8 formats for deep learning,\" arXiv:2209.05433, 2022.",
    "[8]  Z. Dong et al., \"HAWQ: Hessian AWare Quantization of Neural Networks with Mixed-Precision,\" ICCV, 2019.",
    "[9]  V. Mrazek et al., \"ALWANN: Automatic layer-wise approximation for non-iterative accuracy prediction,\" ICCAD, 2019.",
]

txb = slide.shapes.add_textbox(Emu(457200), Emu(838200), Emu(11277295), Emu(5715000))
tf = txb.text_frame
tf.word_wrap = True

for i, ref in enumerate(refs):
    if i == 0:
        para = tf.paragraphs[0]
    else:
        para = tf.add_paragraph()
    para.space_before = Pt(4)
    run = para.add_run()
    run.text = ref
    run.font.size = Pt(10)
    run.font.color.rgb = GRAY
    run.font.bold = False

prs.save(_PPTX)
print('Slide 12 (References) added. Total slides:', len(prs.slides))
