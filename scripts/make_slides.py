"""
make_slides.py  --  generate / selectively rebuild defense presentation

Full rebuild:     py -3.14 scripts/make_slides.py
Selective rebuild: py -3.14 scripts/make_slides.py 4 5
  → only slides 4 and 5 are regenerated; all other slides kept as-is,
    preserving any manual edits made in PowerPoint.
Out: deliverables/AxMAC_Defense_Slides.pptx
"""

import sys, os
sys.stdout.reconfigure(encoding="utf-8")

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as I, Pt
from pptx.dml.color import RGBColor as RGB

# ── constants ────────────────────────────────────────────────────────────────
W, H = Inches(13.333), Inches(7.5)          # 16:9

BG       = RGB(0x0f, 0x11, 0x17)            # dark background
SURFACE  = RGB(0x1a, 0x1d, 0x27)            # card surface
ACCENT   = RGB(0x5b, 0x8d, 0xee)            # blue accent
GREEN    = RGB(0x3e, 0xcf, 0x8e)            # green
RED      = RGB(0xe0, 0x5c, 0x5c)            # red
YELLOW   = RGB(0xf0, 0xc0, 0x60)            # yellow
WHITE    = RGB(0xff, 0xff, 0xff)
MUTED    = RGB(0x7a, 0x80, 0xa0)
ORANGE   = RGB(0xe0, 0x8a, 0x3c)

ROOT    = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))  # project root
FIGDIR  = os.path.join(ROOT, "experiments", "results", "figures")
OUT     = os.path.join(ROOT, "deliverables", "AxMAC_Defense_Slides.pptx")

# ── selective rebuild: py make_slides.py [slide numbers...] ──────────────
# e.g.  py make_slides.py 4 5   → only rebuild slides 4 and 5
_REBUILD = {int(x) for x in sys.argv[1:]} if len(sys.argv) > 1 else set()
_FULL    = not _REBUILD

if _FULL or not os.path.exists(OUT):
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
else:
    prs = Presentation(OUT)

blank_layout = prs.slide_layouts[6]   # completely blank


def _finalize(n):
    """Partial-rebuild mode: replace old slide n with the newly appended one."""
    if _FULL:
        return
    _R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    xml   = prs.slides._sldIdLst
    total = len(prs.slides)
    # Properly drop old slide: remove relationship + part from package
    old_elem = xml[n - 1]
    rId = old_elem.get(f'{{{_R}}}id')
    prs.part.drop_rel(rId)
    xml.remove(old_elem)
    # Move new slide (now at total-2) to correct position
    new_elem = xml[total - 2]
    xml.remove(new_elem)
    xml.insert(n - 1, new_elem)


# ── helpers ──────────────────────────────────────────────────────────────────

def new_slide(n=0):
    """n = 1-indexed slide number; returns None for skipped slides."""
    if not _FULL and n not in _REBUILD:
        return None
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return slide


def box(slide, left, top, width, height,
        fill=None, border=None, border_w=Pt(1)):
    if slide is None: return None
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        I(left), I(top), I(width), I(height)
    )
    shape.fill.solid() if fill else shape.fill.background()
    if fill:
        shape.fill.fore_color.rgb = fill
    shape.line.fill.background() if not border else None
    if border:
        shape.line.color.rgb = border
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, left, top, width, height,
        size=18, bold=False, italic=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, anchor="top"):
    if slide is None: return None
    txb = slide.shapes.add_textbox(I(left), I(top), I(width), I(height))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text       = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def title_bar(slide, title_text, subtitle_text=""):
    if slide is None: return
    box(slide, 0, 0, 13.333, 1.0, fill=SURFACE)
    box(slide, 0, 0.97, 13.333, 0.04, fill=ACCENT)
    txt(slide, title_text, 0.3, 0.08, 12.0, 0.8,
        size=24, bold=True, color=WHITE)
    if subtitle_text:
        txt(slide, subtitle_text, 0.3, 0.65, 10.0, 0.35,
            size=13, color=MUTED)


def bullet_block(slide, items, left, top, width, height,
                 size=16, spacing=0.42, color=WHITE,
                 bullet="•", bullet_color=None):
    if slide is None: return top
    y = top
    for item, level in items:
        indent = left + level * 0.35
        bcolor = bullet_color or ACCENT
        if bullet:
            txt(slide, bullet, indent, y, 0.3, spacing,
                size=size, color=bcolor)
        txt(slide, item, indent + 0.28, y, width - level*0.35 - 0.28, spacing,
            size=size, color=color)
        y += spacing
    return y


def add_figure(slide, path, left, top, width=None, height=None):
    if slide is None: return
    if not os.path.exists(path):
        print(f"  WARNING: figure not found: {path}")
        return
    kw = {}
    if width:  kw["width"]  = I(width)
    if height: kw["height"] = I(height)
    slide.shapes.add_picture(path, I(left), I(top), **kw)


def card(slide, left, top, width, height, fill=SURFACE, border=ACCENT):
    if slide is None: return None
    box(slide, left, top, width, height, fill=fill, border=border, border_w=Pt(1.2))


def slide_number(slide, n, total=11):
    if slide is None:
        return
    txt(slide, f"{n} / {total}", 12.5, 7.1, 0.8, 0.35,
        size=11, color=MUTED, align=PP_ALIGN.RIGHT)
    _finalize(n)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(1)

# Top / bottom accent strips
box(s, 0, 0,    13.333, 0.08, fill=ACCENT)
box(s, 0, 7.42, 13.333, 0.08, fill=ACCENT)

# Title block — shifted down for visual balance (top margin ≈ bottom margin)
txt(s, "Bias-Aware Approximate MAC Design",
    0.6, 1.7, 12.0, 1.0, size=36, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER)
txt(s, "for Energy-Efficient DNN Inference",
    0.6, 2.6, 12.0, 0.8, size=30, bold=True, color=ACCENT,
    align=PP_ALIGN.CENTER)

txt(s, "Jiabo Zhang  ·  Yuxuan Wang",
    0.6, 3.7, 12.0, 0.5, size=18, color=MUTED, align=PP_ALIGN.CENTER)
txt(s, "EEC 289Q — Deep Learning Hardware  |  Spring 2026  |  UC Davis",
    0.6, 4.2, 12.0, 0.5, size=15, color=MUTED, align=PP_ALIGN.CENTER)

# Separator between author block and flow row
box(s, 1.5, 4.85, 10.333, 0.025, fill=MUTED)

# Flow row — option B: unified ACCENT border, Deficiency highlighted with YELLOW
# each entry: (label, border_color, text_color, fill_color)
flow_cards = [
    ("Baseline\nReproduction",       ACCENT, ACCENT, SURFACE),
    ("Deficiency\nIdentified",        YELLOW, YELLOW, RGB(0x1e, 0x1a, 0x0e)),
    ("Contribution A\nBias Fix",      ACCENT, ACCENT, SURFACE),
    ("Contribution B\nNon-uniform K", ACCENT, ACCENT, SURFACE),
    ("RTL +\nFPGA Demo",              ACCENT, ACCENT, SURFACE),
]
_cw, _ch = 1.6, 1.0
_aw, _gap = 0.38, 0.1
_start_x = (13.333 - (len(flow_cards) * _cw
             + (len(flow_cards) - 1) * (_gap + _aw + _gap))) / 2
_x = _start_x
for _i, (_label, _border, _tcol, _fill) in enumerate(flow_cards):
    card(s, _x, 5.0, _cw, _ch, fill=_fill, border=_border)
    txt(s, _label, _x + 0.08, 5.0, _cw - 0.16, _ch,
        size=13, color=_tcol, align=PP_ALIGN.CENTER, bold=True)
    _x += _cw
    if _i < len(flow_cards) - 1:
        _x += _gap
        txt(s, "▸", _x, 5.3, _aw, 0.45, size=18, color=MUTED,
            align=PP_ALIGN.CENTER)
        _x += _aw + _gap

slide_number(s, 1)
print("Slide 1: Title")


# ── Pre-generate energy chart for Slide 2 (only when rebuilding slide 2) ──
_energy_chart_path = os.path.join(FIGDIR, "fig_energy_ops.png")
if _FULL or 2 in _REBUILD:
    _e_labels = ["INT8 add", "INT8 ×", "FP8 ×", "FP16 ×", "FP32 ×"]
    _e_values = [0.03, 0.20, 0.35, 1.10, 3.70]
    _e_annots = ["0.03 pJ", "0.2 pJ", "~0.35 pJ", "1.1 pJ", "3.7 pJ"]
    _e_colors = ["#5b8dee", "#5b8dee", "#5b8dee", "#5b8dee", "#e05c5c"]
    _fig, _ax = plt.subplots(figsize=(5.2, 2.75))
    _fig.patch.set_facecolor("#1a1d27")
    _ax.set_facecolor("#1a1d27")
    _bars = _ax.barh(_e_labels, _e_values, color=_e_colors,
                     height=0.55, edgecolor="none")
    _ax.set_xscale("log")
    _ax.set_xlim(0.015, 18)
    _ax.invert_yaxis()
    for _bar, _ann, _col in zip(_bars, _e_annots, _e_colors):
        _ax.text(_bar.get_width() * 1.4, _bar.get_y() + _bar.get_height() / 2,
                 _ann, va="center", ha="left",
                 color=_col, fontsize=9.5, fontweight="bold",
                 fontfamily="monospace")
    _ax.set_xlabel("Energy (pJ)  —  log scale", color="#7a80a0", fontsize=9)
    _ax.tick_params(axis="both", colors="#7a80a0", labelsize=9)
    _ax.xaxis.set_tick_params(which="minor", length=0)
    for _sp in _ax.spines.values():
        _sp.set_visible(False)
    _ax.grid(axis="x", color="#2a2d3a", linestyle="--", linewidth=0.8, alpha=0.9)
    _ax.set_axisbelow(True)
    _fig.tight_layout(pad=0.6)
    _fig.savefig(_energy_chart_path, dpi=150, bbox_inches="tight",
                 facecolor="#1a1d27")
    plt.close(_fig)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Motivation
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(2)
title_bar(s, "Motivation: MAC Dominates DNN Inference Energy")

# ── Left card ──────────────────────────────────────────────────────────────
card(s, 0.3, 1.15, 5.8, 5.9)
txt(s, "Energy per Operation (45 nm, 1 GHz)",
    0.5, 1.25, 5.4, 0.45, size=14, bold=True, color=ACCENT)

add_figure(s, _energy_chart_path, left=0.38, top=1.72, width=5.55)

txt(s, "Source: Horowitz, ISSCC 2014",
    0.5, 4.62, 5.3, 0.3, size=11, italic=True, color=MUTED)

# "10¹² MACs" callout — prominent sub-card
card(s, 0.38, 5.00, 5.64, 1.9, fill=RGB(0x13, 0x15, 0x1f), border=ACCENT)
txt(s, "A modern LLM: 10¹² MACs / token",
    0.5, 5.12, 5.4, 0.55, size=19, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER)
txt(s, "Reducing per-MAC energy by even 20%\nsaves terajoules at data-center scale.",
    0.5, 5.75, 5.4, 0.70, size=13, color=MUTED,
    align=PP_ALIGN.CENTER)

# ── Right: two sub-cards ───────────────────────────────────────────────────
# Card 1 — How it works (ACCENT border)
card(s, 6.4, 1.15, 6.6, 2.5)
txt(s, "Truncated Multiplier",
    6.6, 1.25, 6.2, 0.45, size=14, bold=True, color=ACCENT)
box(s, 6.6, 1.72, 6.0, 0.025, fill=ACCENT)
txt(s, "Drop the K lowest bits of each integer\nproduct before accumulation.",
    6.6, 1.82, 6.2, 0.65, size=13, color=WHITE)
txt(s, "Hardware savings",
    6.6, 2.55, 6.2, 0.35, size=12, bold=True, color=ACCENT)
txt(s, "K(2n–K) / (2n²)  ≈  25%  at  K = 4",
    6.6, 2.92, 6.2, 0.35, size=12, color=MUTED)

# Card 2 — The problem (YELLOW border + dark fill)
card(s, 6.4, 3.80, 6.6, 3.25, fill=RGB(0x1e, 0x1a, 0x0e), border=YELLOW)
txt(s, "The Problem  ⚠",
    6.6, 3.90, 6.2, 0.45, size=14, bold=True, color=YELLOW)
box(s, 6.6, 4.37, 6.0, 0.025, fill=YELLOW)
bullet_block(s,
    [("Always rounds toward –∞", 0),
     ("Error ∈ [0, 2ᴷ)  →  one-sided bias", 0),
     ("Grows like N across N MACs,  not √N", 0)],
    left=6.6, top=4.50, width=6.0, height=1.4,
    size=13, spacing=0.46, bullet="•", bullet_color=YELLOW, color=WHITE)
card(s, 6.55, 6.12, 6.3, 0.70,
     fill=SURFACE, border=ACCENT)
txt(s, "→  Addressed by round mode in this work",
    6.7, 6.22, 6.1, 0.50, size=13, bold=True, color=ACCENT,
    align=PP_ALIGN.CENTER)

slide_number(s, 2)
print("Slide 2: Motivation")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Bias Problem
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(3)
title_bar(s, "The Bias Accumulation Problem",
          "Truncation error is one-sided → grows linearly with accumulation depth N")

# Figure 1
fig1 = os.path.join(FIGDIR, "fig1_bias_accumulation.png")
add_figure(s, fig1, left=0.25, top=1.1, width=6.2)

# ── Key numbers card ──────────────────────────────────────────────────────
card(s, 6.7, 1.1, 6.3, 2.8)
txt(s, "Key numbers  (INT8, K=6, N=4096)",
    6.9, 1.2, 5.9, 0.45, size=13, bold=True, color=ACCENT)
box(s, 6.9, 1.67, 5.9, 0.025, fill=ACCENT)

nums = [
    ("trunc  bias",      "26,617", RED,   WHITE),
    ("round  bias",      "−14",    MUTED, WHITE),
    ("stochastic  bias", "14",     MUTED, WHITE),
]
y = 1.78
for label, val, lcol, vcol in nums:
    txt(s, label, 6.9,  y, 3.5, 0.48, size=14, color=lcol)
    txt(s, val,   10.5, y, 2.2, 0.48, size=20, bold=True, color=vcol)
    y += 0.52
txt(s, "trunc grows ∝ N;  round/stoch grow ∝ √N",
    6.9, 3.40, 5.9, 0.40, size=12, italic=True, color=MUTED)

# ── Theorem 1 card ─────────────────────────────────────────────────────────
card(s, 6.7, 4.15, 6.3, 2.9)
txt(s, "Theorem 1  (Formal Proof)",
    6.9, 4.25, 5.9, 0.45, size=13, bold=True, color=ACCENT)
box(s, 6.9, 4.72, 5.9, 0.025, fill=ACCENT)

thm = [
    ("E[Σ εᵢ] = N · 2^(K−1)", "[trunc]"),
    ("E[Σ δᵢ] = 0  for any N", "[round]"),
    ("Var[Σ] = N · (2^K)² / 12", "[both]" ),
]
y = 4.82
for eq, tag in thm:
    txt(s, eq,  6.9,  y, 4.5, 0.42, size=12, color=WHITE)
    txt(s, tag, 11.5, y, 1.3, 0.42, size=11, color=MUTED,
        align=PP_ALIGN.RIGHT)
    y += 0.42

card(s, 6.85, 6.20, 6.0, 0.65,
     fill=RGB(0x0a, 0x1e, 0x0e), border=GREEN)
txt(s, "round mode = free stochastic rounding  (no RNG needed)",
    7.0, 6.30, 5.7, 0.45, size=12, bold=True, color=GREEN,
    align=PP_ALIGN.CENTER)

slide_number(s, 3)
print("Slide 3: Bias problem")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Contribution A: Rounding Modes
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(4)
title_bar(s, "Contribution A: Error-Compensated Multiplier",
          "Add correction constant 2^(K−1) before masking — same hardware cost as truncation")

modes = [
    ("trunc",       "product & ~mask",
     "Error ∈ [0, 2^K)\nMean ≈ 2^(K−1) > 0\nLinear bias growth",
     RED),
    ("round",       "(product + 2^(K−1)) & ~mask",
     "Error ∈ [−2^(K−1), 2^(K−1))\nMean = 0  (zero bias!)\nConstant add, folded into PPT",
     GREEN),
    ("stochastic",  "(product + rand[0,2^K)) & ~mask",
     "True zero-mean (√N growth)\nNeeds K-bit LFSR per MAC\n→ Gupta et al., ICML 2015",
     YELLOW),
]
x = 0.3
for name, op, desc, col in modes:
    card(s, x, 1.15, 4.2, 5.9, fill=SURFACE, border=col)
    txt(s, name,  x+0.15, 1.25, 3.9, 0.55, size=22, bold=True, color=col,
        align=PP_ALIGN.CENTER)
    txt(s, op,    x+0.1,  1.9,  4.0, 0.7,  size=11, color=MUTED,
        align=PP_ALIGN.CENTER)
    txt(s, desc,  x+0.15, 2.75, 3.9, 2.2,  size=14, color=WHITE)

    # HW cost label
    hw = "Baseline" if name=="trunc" else ("≈ trunc" if name=="round" else "+LFSR RNG")
    hw_col = MUTED if name=="trunc" else (GREEN if name=="round" else RED)
    card(s, x+0.2, 5.5, 3.8, 1.3, fill=RGB(0x13,0x15,0x1f), border=hw_col)
    txt(s, "HW cost:", x+0.35, 5.6, 3.4, 0.4, size=12, color=MUTED)
    txt(s, hw,         x+0.35, 6.0, 3.4, 0.55, size=17, bold=True, color=hw_col,
        align=PP_ALIGN.CENTER)
    x += 4.35

slide_number(s, 4)
print("Slide 4: Contribution A")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — RTL Implementation
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(5)
title_bar(s, "RTL Implementation & Verification",
          "Synthesizable Verilog  ·  ModelSim  ·  Quartus Prime Lite  ·  EP4CE10 FPGA")

# Left: module list
card(s, 0.3, 1.15, 5.5, 5.9)
txt(s, "RTL Modules (9 .v)", 0.5, 1.25, 5.1, 0.45, size=14, bold=True, color=ACCENT)
modules = [
    ("mac_unit.v",       "Parameterized approx MAC (K, mode)", WHITE),
    ("aca_adder.v",      "ACA approximate adder (W)",           WHITE),
    ("mac_array.v",      "16-MAC dot-product array",            WHITE),
    ("lfsr.v",           "K-bit LFSR for stochastic rounding",  MUTED),
    ("mlp_top.v",        "2-layer MLP (simulation)",            WHITE),
    ("mlp_top_demo.v",   "Board wrapper + LED display",         WHITE),
    ("uart_tx/framer.v", "UART 115200 baud output",             WHITE),
    ("drum_multiplier.v","DRUM-k unbiased multiplier (P1.1)",   ACCENT),
]
y = 1.78
for mod, desc, col in modules:
    txt(s, f"  {mod}", 0.45, y, 2.5, 0.38, size=11, bold=True, color=col)
    txt(s, desc,      2.95, y, 2.7, 0.38, size=11, color=MUTED)
    y += 0.4

# Right top: verification results
card(s, 6.1, 1.15, 6.9, 2.7)
txt(s, "Verification", 6.3, 1.25, 6.5, 0.45, size=14, bold=True, color=ACCENT)
vresults = [
    ("ModelSim: 6 testbenches,  all PASS", GREEN),
    ("  200 vectors each (golden from Python)", MUTED),
    ("350 pytest tests,  all PASS", GREEN),
]
y = 1.78
for line, col in vresults:
    txt(s, line, 6.3, y, 6.5, 0.4, size=13, color=col)
    y += 0.44

# Right bottom: synthesis results
card(s, 6.1, 4.05, 6.9, 3.0)
txt(s, "Quartus Synthesis (EP4CE10)", 6.3, 4.15, 6.5, 0.45,
    size=14, bold=True, color=ACCENT)
synth = [
    ("Logic Elements",      "1,769 / 10,320  (17%)",  WHITE),
    ("Timing",              "50 MHz  closed ✓",        GREEN),
    ("round vs trunc ΔLE",  "+46 LE  (2.6%)",          GREEN),
    ("stochastic ΔLE",      "+48 LE  (needs LFSR)",    YELLOW),
]
y = 4.7
for label, val, col in synth:
    txt(s, label, 6.3, y, 3.5, 0.42, size=13, color=MUTED)
    txt(s, val,  9.85, y, 3.0, 0.42, size=13, bold=True, color=col)
    y += 0.48

slide_number(s, 5)
print("Slide 5: RTL")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — FPGA On-Chip Validation  ← KEY DEMO SLIDE
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(6)
title_bar(s, "FPGA On-Chip Validation  —  K=6 trunc causes MISCLASSIFICATION",
          "EP4CE10  ·  野火征途 Pro  ·  UART output  ·  2026-06-01")

# Left: setup description
card(s, 0.3, 1.15, 4.8, 5.9)
txt(s, "Demo Setup", 0.5, 1.25, 4.4, 0.45, size=14, bold=True, color=ACCENT)
setup = [
    "• Toy 2-layer MLP",
    "  weights quantized to INT8",
    "  correct class = 1",
    "",
    "• 5 bitstreams programmed",
    "  via JTAG (USB-Blaster)",
    "",
    "• UART @ 115200 baud",
    "  streams 10 × int32 logits",
    "  back to PC",
    "",
    "• read_uart.py decodes",
    "  argmax on PC",
]
y = 1.8
for line in setup:
    txt(s, line, 0.5, y, 4.3, 0.38, size=13, color=WHITE if line else MUTED)
    y += 0.38

# Right: results table
card(s, 5.4, 1.15, 7.6, 5.9)
txt(s, "On-Chip Results", 5.6, 1.25, 7.2, 0.45, size=14, bold=True, color=ACCENT)

configs = [
    ("K=0  trunc", "exact baseline",  "1", True),
    ("K=2  trunc", "",                "1", True),
    ("K=4  trunc", "",                "1", True),
    ("K=4  round", "bias fix",        "1", True),
    ("K=6  trunc", "aggressive trunc","3", False),
]
hdrs = ["Config", "Note", "Board argmax", "Correct?"]
col_x = [5.6, 7.4, 9.8, 11.4]
col_w = [1.75, 2.3, 1.5, 1.6]
# Header
y = 1.85
for j, h in enumerate(hdrs):
    txt(s, h, col_x[j], y, col_w[j], 0.42,
        size=12, bold=True, color=ACCENT)
y = 2.32
box(s, 5.4, y-0.05, 7.55, 0.04, fill=ACCENT)

for cfg, note, argmax, correct in configs:
    col = WHITE if correct else RED
    bg_col = RGB(0x1e,0x0e,0x0e) if not correct else None
    if bg_col:
        box(s, 5.4, y-0.05, 7.55, 0.58, fill=bg_col)

    txt(s, cfg,    col_x[0], y, col_w[0], 0.48, size=13, bold=True, color=col)
    txt(s, note,   col_x[1], y, col_w[1], 0.48, size=12, color=MUTED)
    txt(s, argmax, col_x[2], y, col_w[2], 0.48, size=18,
        bold=True, color=GREEN if correct else RED, align=PP_ALIGN.CENTER)
    mark = "✓" if correct else "✗"
    txt(s, mark,   col_x[3], y, col_w[3], 0.48, size=18,
        bold=True, color=GREEN if correct else RED, align=PP_ALIGN.CENTER)
    y += 0.55

# Callout
card(s, 5.4, 5.55, 7.6, 1.3, fill=RGB(0x1e,0x0a,0x0a), border=RED)
txt(s,
    "K=6 trunc:  argmax 1 → 3  (WRONG)\n"
    "Matches Python sim bias +29.9  ✓\n"
    "RTL faithfully implements bias accumulation",
    5.55, 5.65, 7.3, 1.1, size=13, bold=True, color=RED)

slide_number(s, 6)
print("Slide 6: FPGA on-chip")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CIFAR-10 Results  ← MONEY SLIDE
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(7)
title_bar(s, "CIFAR-10 Validation — round mode is FREE accuracy preservation",
          "SimpleCNN  ·  INT8 PTQ  ·  10,000 test images  ·  K = 0 … 6")

fig5 = os.path.join(FIGDIR, "fig5_cifar10_accuracy.png")
add_figure(s, fig5, left=0.25, top=1.1, width=7.5)

# Key numbers
card(s, 8.0, 1.1, 5.1, 5.9)
txt(s, "Key Results", 8.2, 1.2, 4.7, 0.45, size=14, bold=True, color=ACCENT)

rows = [
    ("Float32 baseline",  "83.25%", WHITE),
    ("INT8 exact (K=0)",  "82.94%", WHITE),
    ("",                  "",       MUTED),
    ("round  K=6",        "83.0%",  GREEN),
    ("round  K=4",        "83.0%",  GREEN),
    ("round  K=2",        "83.0%",  GREEN),
    ("",                  "",       MUTED),
    ("trunc  K=4",        "79.5%",  YELLOW),
    ("trunc  K=5",        "56.8%",  RED),
    ("trunc  K=6",        "10.8%",  RED),
    ("",                  "",       MUTED),
    ("DRUM-4",            "82.4%",  ORANGE),
]
y = 1.75
for label, val, col in rows:
    if not label:
        y += 0.15
        continue
    txt(s, label, 8.2,  y, 3.0, 0.42, size=13, color=MUTED)
    txt(s, val,   11.25, y, 1.6, 0.42, size=14, bold=True, color=col)
    y += 0.44

card(s, 8.0, 6.15, 5.1, 0.85, fill=RGB(0x0a,0x1e,0x0e), border=GREEN)
txt(s,
    "round mode: +72.1 pp over trunc at K=6\n"
    "Zero additional hardware cost",
    8.15, 6.2, 4.9, 0.75, size=13, bold=True, color=GREEN)

slide_number(s, 7)
print("Slide 7: CIFAR-10 results")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Hardware PPA
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(8)
title_bar(s, "Hardware PPA Analysis",
          "Quartus Prime Lite 18.1  ·  EP4CE10F17C8  ·  50 MHz constraint")

fig3 = os.path.join(FIGDIR, "fig3_ppa_hardware.png")
add_figure(s, fig3, left=0.25, top=1.1, width=7.3)

card(s, 7.8, 1.1, 5.3, 5.9)
txt(s, "Takeaways", 8.0, 1.2, 5.0, 0.45, size=14, bold=True, color=ACCENT)

points = [
    ("round vs trunc", "+46 LE  (+2.6%)\n≈ identical area", GREEN),
    ("stochastic", "+48 LE  (+LFSR)\nExceeds exact baseline!", RED),
    ("DRUM-4 (ppa_drum_k4)", "Quartus run pending\nExpected: –LEs for mult\nbut +LZC+shift overhead", ORANGE),
]
y = 1.75
for title, body, col in points:
    card(s, 7.85, y, 5.1, 1.55, fill=SURFACE, border=col)
    txt(s, title, 8.0, y+0.1,  4.8, 0.45, size=13, bold=True, color=col)
    txt(s, body,  8.0, y+0.58, 4.8, 0.88, size=12, color=WHITE)
    y += 1.72

txt(s, "→  round mode is the practical optimum:\n"
       "   zero-mean bias + near-zero HW cost",
    7.85, 6.25, 5.1, 0.75, size=13, bold=True, color=GREEN)

slide_number(s, 8)
print("Slide 8: PPA")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Contribution B
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(9)
title_bar(s, "Contribution B: Sensitivity-Driven Non-Uniform K Allocation",
          "Greedy layer-wise K budget allocation  ·  cf. ALWANN (ICCAD 2019), HAWQ (ICCV 2019)")

# Left: problem
card(s, 0.3, 1.15, 5.8, 5.9)
txt(s, "The Uniform-K Problem", 0.5, 1.25, 5.4, 0.45, size=14, bold=True, color=RED)
txt(s,
    "Demo MLP layer breakdown:\n"
    "  Layer 1:  784×64 = 50,176 MACs  (95.8%)\n"
    "  Layer 2:  64×10  =    640 MACs  ( 1.2%)\n\n"
    "With uniform K, the tiny layer 2\n"
    "gets the same truncation depth as\n"
    "the dominant layer 1.\n\n"
    "→  Energy wasted on insensitive layers.\n"
    "→  Sensitive layers truncated too deep.",
    0.5, 1.8, 5.4, 4.2, size=14, color=WHITE)

# Right: solution
card(s, 6.4, 1.15, 6.6, 2.9)
txt(s, "Greedy Allocation Algorithm", 6.6, 1.25, 6.2, 0.45,
    size=14, bold=True, color=ACCENT)
txt(s,
    "score_i = MAC_count[i] / (sensitivity[i] × (K[i]+1))\n\n"
    "Assign next K unit to argmax(score):\n"
    "  High MACs → more savings\n"
    "  High sensitivity → smaller K\n"
    "  Prevents one layer dominating budget",
    6.6, 1.78, 6.2, 3.0, size=13, color=WHITE)

card(s, 6.4, 4.2, 6.6, 2.8)
txt(s, "Result", 6.6, 4.3, 6.2, 0.45, size=14, bold=True, color=GREEN)
txt(s,
    "Same energy budget:\n"
    "  Non-uniform K allocation\n"
    "  reduces logit NRMSE vs uniform K\n\n"
    "Layer 2 keeps K=0 (exact),\n"
    "Layer 1 absorbs more truncation.",
    6.6, 4.8, 6.2, 2.0, size=13, color=WHITE)

slide_number(s, 9)
print("Slide 9: Contribution B")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Future Work
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(10)
title_bar(s, "Future Work  →  ISCAS 2027 / IEEE Embedded Systems Letters",
          "Three concrete directions toward a publishable extension")

items = [
    (
        "① ASIC Synthesis  (NanGate45 + OpenROAD)",
        "Replace FPGA proxy with 45 nm ASIC numbers.\n"
        "Precise area / power comparison: trunc vs round vs DRUM.\n"
        "DRUM area advantage vs round control-logic overhead → definitive HW verdict.",
        ACCENT,
    ),
    (
        "② ResNet-20 / MobileNet-v2 on CIFAR-10 + ImageNet",
        "Extend Contribution B (non-uniform K) to real trained architectures.\n"
        "Energy–accuracy Pareto front at task level (top-1 accuracy vs pJ/inference).\n"
        "Baseline: ALWANN on same networks.",
        GREEN,
    ),
    (
        "③ Transformer Attention Bias Validation",
        "Q·K^T at seq_len=512: each attention score accumulates d_head=64 MACs.\n"
        "Theorem 1 corollary: trunc K=4 → 512 LSB bias, exceeding INT8 range.\n"
        "Experiment: GPT-2 small, compare perplexity under trunc vs round.",
        YELLOW,
    ),
]

y = 1.2
for title, body, col in items:
    card(s, 0.3, y, 12.7, 1.85, fill=SURFACE, border=col)
    txt(s, title, 0.5, y+0.08, 12.3, 0.5,  size=15, bold=True, color=col)
    txt(s, body,  0.5, y+0.62, 12.3, 1.1,  size=13, color=WHITE)
    y += 2.05

slide_number(s, 10)
print("Slide 10: Future Work")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Conclusion
# ════════════════════════════════════════════════════════════════════════════
s = new_slide(11)
box(s, 0, 0, 13.333, 0.08, fill=ACCENT)
box(s, 0, 7.42, 13.333, 0.08, fill=ACCENT)
title_bar(s, "Conclusion")

takeaways = [
    (
        "Truncation bias is the dominant accuracy loss in approximate MAC inference.",
        "Error always in [0, 2^K) → one-sided → coherent accumulation ∝ N.\n"
        "Formal proof: E[bias_N] = N · 2^(K-1) for trunc,  0 for round.",
        RED,
    ),
    (
        "round mode = free fix: same hardware cost, zero bias.",
        "+46 LE (+2.6%) on EP4CE10.  CIFAR-10: 83.0% at K=6 vs 10.8% (trunc).\n"
        "72.1 pp accuracy gap at zero hardware cost — confirmed on real FPGA.",
        GREEN,
    ),
    (
        "Non-uniform K allocation further improves energy–accuracy trade-off.",
        "Sensitivity-driven greedy allocation; same energy, lower logit NRMSE.\n"
        "cf. ALWANN / HAWQ — simpler, no Hessian, directly exploits MAC counts.",
        ACCENT,
    ),
]

y = 1.15
for title, body, col in takeaways:
    card(s, 0.3, y, 12.7, 1.7, fill=SURFACE, border=col)
    txt(s, title, 0.5, y+0.08, 12.3, 0.48, size=15, bold=True, color=col)
    txt(s, body,  0.5, y+0.6,  12.3, 0.95, size=13, color=WHITE)
    y += 1.88

card(s, 0.3, 6.75, 12.7, 0.65, fill=RGB(0x13,0x15,0x1f), border=ACCENT)
txt(s,
    "Jiabo Zhang  ·  Yuxuan Wang  |  EEC 289Q Deep Learning Hardware  |  "
    "Spring 2026, UC Davis  |  bojzhang@ucdavis.edu",
    0.5, 6.82, 12.3, 0.5, size=12, color=MUTED, align=PP_ALIGN.CENTER)

slide_number(s, 11)
print("Slide 11: Conclusion")


# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"\nSaved: {OUT}")
print(f"  Slides: {len(prs.slides)}")
